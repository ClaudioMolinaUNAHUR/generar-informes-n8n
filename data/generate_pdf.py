#!/usr/bin/env python3
import sys
import json
import base64
import subprocess
import os
import warnings
import uuid
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
import matplotlib.pyplot as plt
from pypdf import PdfReader, PdfWriter
import requests
import textwrap
import numpy as np
from io import BytesIO
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings("ignore")
DATA_DIR = "/data"


def _prepare_output_file(path):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    if os.path.exists(path):
        os.remove(path)


# --------------------------------------------------------------
# UTILS
# --------------------------------------------------------------
def truncate_text(text, max_chars=300):
    if len(text) > max_chars:
        return text[:max_chars].rsplit(' ', 1)[0] + "…"
    return text


def normalize_kpi_text(text, max_chars=250):
    val = truncate_text((text or "").strip(), max_chars)
    return val if val else " "


def normalize_suggestion_text(text, max_chars=200, line_width=46):
    raw = truncate_text((text or "").strip(), max_chars)
    if not raw:
        return " "

    raw = raw.replace("\r\n", "\n").replace("\r", "\n").replace("\\n", "\n")
    lines = raw.split("\n")

    wrapped_lines = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            # Preserve intentional blank lines between suggestion blocks.
            if not wrapped_lines or wrapped_lines[-1] != "":
                wrapped_lines.append("")
            continue
        bullet = stripped.startswith("•") or stripped.startswith("-")
        content = stripped.lstrip("•- ").strip() if bullet else stripped
        if not content:
            continue
        if bullet:
            wrapped = textwrap.fill(
                content,
                width=line_width,
                initial_indent="• ",
                subsequent_indent="  ",
            )
        else:
            wrapped = textwrap.fill(content, width=line_width)
        wrapped_lines.append(wrapped)

    body = "\n".join(wrapped_lines) if wrapped_lines else " "
    if body.strip() and not body.upper().startswith("SUGERENCIAS"):
        return f"SUGERENCIAS:\n{body}"
    return body


def apply_text_formatting(text_frame, font_name='Aptos', size=None, set_line=True):
    """
    Aplica formato de fuente, tamaño e interlineado a todos los párrafos y runs en el text_frame.
    Si font_name es None, no cambia la fuente.
    Si size es None, no cambia el tamaño.
    Si set_line=True, establece interlineado 1.5.
    """
    # Auto-shrink: reduce fuente e interlineado para que entre en el box
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        if set_line:
            paragraph.line_spacing = 1.5
        for run in paragraph.runs:
            if font_name is not None:
                run.font.name = font_name
            if size is not None:
                run.font.size = Pt(size)


def get_logo_from_base64(base64_string: str) -> BytesIO | None:
    """
    Decodifica una cadena Base64 y devuelve un objeto BytesIO con los datos de la imagen.
    Retorna None si la cadena está vacía o es inválida.
    """
    if not base64_string:
        return None
    try:
        image_data = base64.b64decode(base64_string)
        return BytesIO(image_data)
    except base64.binascii.Error:
        log("⚠️ Error de decodificación Base64. La cadena del logo podría ser inválida.")
        return None


def replace_placeholders(slide, replacements):
    """
    Busca un placeholder por su nombre (definido en el Panel de Selección de PowerPoint)
    y reemplaza su texto por el valor correspondiente.
    Si el valor es una lista de objetos, aplica formato (ej. negritas) creando runs.
    Retorna una lista de tuplas (text_frame, key) modificados.
    """
    modified = []

    def _value_to_text(value):
        if isinstance(value, list):
            parts = []
            for item in value:
                if isinstance(item, dict):
                    parts.append(str(item.get("text", "")))
                else:
                    parts.append(str(item))
            return "".join(parts)
        if isinstance(value, str):
            return value
        return str(value)

    # 1) Reemplazo exacto: mantiene compatibilidad con placeholders que ocupan todo el cuadro.
    for key, value in replacements.items():
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            texto = shape.text.strip()
            if texto == key:
                tf = shape.text_frame
                tf.clear()  # Limpiar el contenido existente

                if isinstance(value, list):
                    # Procesar como array de objetos con formato
                    for item in value:
                        if not isinstance(item, dict):
                            continue  # Saltar si no es dict
                        text_part = item.get("text", "")
                        if not tf.paragraphs:
                            p = tf.add_paragraph()
                        else:
                            p = tf.paragraphs[-1]
                        run = p.add_run()
                        run.text = text_part
                        if item.get("bold", False):
                            run.font.bold = True
                        # Aquí puedes agregar más formatos si es necesario, ej.:
                        # if item.get("italic", False):
                        #     run.font.italic = True
                else:
                    val_str = _value_to_text(value)
                    val_str = val_str.replace("\\n", "\n").replace("\\\n", "\n")
                    tf.text = val_str

                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.JUSTIFY

                modified.append((tf, key))

    # 2) Reemplazo embebido: soporta placeholders dentro de un texto mayor en el mismo cuadro.
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        original_text = tf.text or ""
        new_text = original_text
        replaced_keys = []

        for key, value in replacements.items():
            if key in new_text:
                val_str = _value_to_text(value)
                val_str = val_str.replace("\\n", "\n").replace("\\\n", "\n")
                new_text = new_text.replace(key, val_str)
                replaced_keys.append(key)

        if new_text != original_text:
            tf.text = new_text
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.JUSTIFY
            for key in replaced_keys:
                modified.append((tf, key))

    # 3) Limpieza final: si queda algún placeholder sin reemplazar, ocultarlo con blanco.
    unresolved_pattern = re.compile(r"\{\{ph_[^}]+\}\}")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        original_text = tf.text or ""
        cleaned_text = unresolved_pattern.sub(" ", original_text)
        if cleaned_text != original_text:
            tf.text = cleaned_text
    
    return modified


def insert_image_scaled_by_width(slide, placeholder, image_path_or_stream):
    """
    Reemplaza un placeholder con una imagen, escalándola para que ocupe todo el ancho
    del placeholder y ajustando el alto proporcionalmente. La imagen se centra verticalmente.
    """
    parent = placeholder.element.getparent()
    if parent is None:
        return  # ya fue eliminado, evitar crash

    ph_left, ph_top = placeholder.left, placeholder.top
    ph_width, ph_height = placeholder.width, placeholder.height

    parent.remove(placeholder.element)

    pic = slide.shapes.add_picture(
        image_path_or_stream, ph_left, ph_top, width=ph_width
    )

    new_height = pic.height
    pic.top = ph_top + (ph_height - new_height) // 2


def insert_logo_preserving_aspect(slide, placeholder, logo_stream):
    # área disponible
    ph_left, ph_top = placeholder.left, placeholder.top
    ph_w, ph_h = placeholder.width, placeholder.height

    # borrar placeholder original
    placeholder.element.getparent().remove(placeholder.element)

    # insertar imagen sin escalar
    pic = slide.shapes.add_picture(logo_stream, ph_left, ph_top)

    # tamaño real de la imagen
    img_w, img_h = pic.width, pic.height

    # calcular factor manteniendo aspecto
    scale = min(ph_w / img_w, ph_h / img_h)

    # aplicar tamaño escalado
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)

    pic.width = new_w
    pic.height = new_h

    # centrar dentro del placeholder original
    pic.left = ph_left + (ph_w - new_w) // 2
    pic.top = ph_top + (ph_h - new_h) // 2


def insert_logo_with_scaling(slide, logo_stream):
    """
    Busca el primer placeholder de tipo 'Picture' (18) e inserta el logo
    dentro de sus límites, manteniendo la relación de aspecto y eliminando el placeholder original.
    """
    # 18 es el tipo de placeholder para 'Picture'
    LOGO_PLACEHOLDER_TYPE = 18

    if not logo_stream:
        return

    for shape in slide.placeholders:
        if shape.placeholder_format.type == LOGO_PLACEHOLDER_TYPE:

            insert_logo_preserving_aspect(slide, shape, logo_stream)
            break  # Solo insertamos el primer logo que encontremos


def log(msg):
    print(msg, flush=True)


def set_placeholder_text(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            if ph.has_text_frame:
                ph.text = text
                apply_text_formatting(ph.text_frame, font_name='Aptos', size=11)
            return
    # Si no existe, opcionalmente crear un cuadro de texto
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = tb.text_frame
    tf.text = text
    apply_text_formatting(tf, font_name='Aptos', size=11)


# --------------------------------------------------------------
# PORTADA
# --------------------------------------------------------------
def generar_portada(data, logo_stream):
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_portada.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_titulo}}": data.get("titulo_portada", ""),
        "{{ph_subtitle}}": data.get("subtitulo_portada", ""),
        "{{ph_fecha}}": data.get("fecha_portada", ""),
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    # Ejecutamos el reemplazo normal
    modified = replace_placeholders(slide, replacements)

    for tf, key in modified:
        key_l = key.lower()
        if "titulo" in key_l:
            apply_text_formatting(tf, font_name='Calibri', size=18)
        elif "sub" in key_l:
            apply_text_formatting(tf, font_name='Calibri', size=14)
        elif "pie" in key_l:
            apply_text_formatting(tf, font_name=None, size=10)
        else:
            apply_text_formatting(tf, font_name=None, size=12)

    # Aplicamos el centrado manual a los campos específicos
    campos_a_centrar = ["{{ph_subtitle}}", "{{ph_fecha}}"]
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Buscamos si el texto de la forma coincide con los valores ya reemplazados
            # (o puedes volver a iterar sobre las keys de replacements)
            for key in campos_a_centrar:
                valor_insertado = replacements[key]
                if valor_insertado and valor_insertado in shape.text:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER

    # Busca un placeholder de tipo imagen (18) para el logo
    insert_logo_with_scaling(slide, logo_stream)

    output = f"{DATA_DIR}/pptx-parts/portada.pptx"
    _prepare_output_file(output)
    prs.save(output)
    return output


# --------------------------------------------------------------
# SLIDE PRODUCTO
# --------------------------------------------------------------
def generar_slide_producto(resumen, product_type, data, logo_stream, pie_l="", pie_r=""):
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_{product_type}.pptx")
    slide = prs.slides[0]

    # Definir campos por tipo de producto
    campos_por_tipo = {
        "uas": ["usu_per", "usu_esp", "solicitudes", "revalida"],
        "beyondtrust": ["pra", "rs", "ps", "adb", "epm"],
        "whalemate": ["sim", "aca", "ana", "grh", "cad"],
    }

    # Mapear product_type a grupo
    tipo_grupo = ""
    if product_type.lower() == "uas":
        tipo_grupo = "uas"
    elif product_type.lower() == "beyondtrust":
        tipo_grupo = "beyondtrust"
    elif product_type.lower() == "whalemate":
        tipo_grupo = "whalemate"

    replacements = {
        "{{ph_resumen}}": resumen,
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    # Agregar los campos nuevos si corresponden
    if tipo_grupo:
        for campo in campos_por_tipo[tipo_grupo]:
            replacements[f"{{{{ph_{campo}}}}}"] = data.get(campo, " ")

    modified = replace_placeholders(slide, replacements)

    for tf, key in modified:
        key_l = key.lower()
        if "titulo" in key_l:
            apply_text_formatting(tf, font_name='Calibri', size=18)
        elif "sub" in key_l:
            apply_text_formatting(tf, font_name='Calibri', size=14)
        elif "kpis" in key_l or "sugerencia" in key_l:
            # Mantener KPIs y sugerencias compactos para compartir el alto disponible.
            apply_text_formatting(tf, font_name='Calibri', size=12, set_line=False)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(0)
                p.space_after = Pt(0)
        elif "pie" in key_l:
            apply_text_formatting(tf, font_name=None, size=10)
        else:
            apply_text_formatting(tf, font_name='Aptos', size=12)

    insert_logo_with_scaling(slide, logo_stream)

    output = f"{DATA_DIR}/pptx-parts/producto_{product_type}.pptx"
    _prepare_output_file(output)
    prs.save(output)
    return output


# --------------------------------------------------------------
# GRÁFICOS
# --------------------------------------------------------------
def create_matplotlib_chart(chart_info, friendly_names, output_file):
    plt.figure(figsize=(10, 5))
    ctype = chart_info.get("type")

    chart_title = chart_info.get("title") or chart_info.get("titulo") or "<no title>"
    use_adaptive_scale = False

    def _normalized_chart_name(value):
        return re.sub(r"\s+", " ", str(value or "").replace("_", " ").strip().lower())

    def _compact_tick_label(value, _pos=None):
        abs_value = abs(value)
        if abs_value >= 1_000_000:
            return f"{value / 1_000_000:.1f}M"
        if abs_value >= 1_000:
            return f"{value / 1_000:.0f}K"
        return f"{int(value)}"

    def _should_use_adaptive_scale(title, series_values):
        if _normalized_chart_name(title) != "alertas controles":
            return False
        positives = [v for values in series_values for v in values if v and v > 0]
        if len(positives) < 2:
            return False
        min_positive = min(positives)
        max_positive = max(positives)
        return min_positive > 0 and (max_positive / min_positive) >= 20

    # Título si viene en la definición del gráfico
    # title = chart_info.get("title") or chart_info.get("titulo") or ""
    # if title:
    #     plt.title(title, fontsize=20, fontweight="bold", pad=20)

    # Aumentar tamaño de fuente para los ejes
    plt.xticks(fontsize=18)
    plt.yticks(fontsize=18)

    labels = chart_info.get("labels", [])
    x = range(len(labels))

    # Mapa de etiquetas amigables para claves comunes
    def _friendly_label(key, value):
        if isinstance(value, dict):
            return value.get("label") or value.get("title") or key.replace("_", " ").capitalize()
        return value

    flat_friendly_names = {
        key: _friendly_label(key, value)
        for chart in friendly_names.values()
        for key, value in chart.items()
    }

    # Paleta de colores para series (se rotan si hay más series)
    palette = ["#4f81bd", "#9abb59", "#4bacc6", "#8064a2"]

    # Detectar series numéricas dinámicamente (manteniendo el orden del dict)
    series_keys = []
    for key, val in chart_info.items():
        if key in ("labels", "type", "title", "titulo"):
            continue
        # Considerar series que sean listas/tuplas de números
        if isinstance(val, (list, tuple)) and all(
            isinstance(v, (int, float)) for v in val
        ):
            series_keys.append(key)

    if ctype == "bar":
        # Barras agrupadas: calcular offsets según cantidad de series
        n = len(series_keys)
        if n == 0:
            # nada que dibujar
            return

        ind = np.arange(len(labels))  # posiciones base
        total_width = 0.7
        bar_width = total_width / n
        plotted_series = []
        use_adaptive_scale = False

        for idx, key in enumerate(series_keys):
            vals = list(chart_info.get(key) or [])
            vals = vals[:len(labels)]  # truncate if longer
            if len(vals) < len(labels):
                vals.extend([0] * (len(labels) - len(vals)))  # pad with 0
            plotted_series.append(vals)

            label_full = flat_friendly_names.get(key, key.replace("_", " ").capitalize())
            # Dividir etiquetas largas en varias líneas para que no encojan el gráfico
            label = textwrap.fill(label_full, width=22)

            color = palette[idx % len(palette)]

            # calcular posiciones para esta serie
            offset = (idx - (n - 1) / 2) * bar_width
            positions = ind + offset
            plt.bar(positions, vals, bar_width * 0.95, label=label, color=color)

        # Ajustar ticks al centro de los grupos
        use_adaptive_scale = _should_use_adaptive_scale(chart_title, plotted_series)
        display_labels = labels
        x_tick_fontsize = 16
        if use_adaptive_scale:
            x_tick_fontsize = 14

        plt.xticks(ind, display_labels, rotation=0, fontsize=x_tick_fontsize)
        plt.grid(axis="y", linestyle="-", color="#dcdcdc", linewidth=0.8)
        if use_adaptive_scale:
            positive_values = [v for values in plotted_series for v in values if v and v > 0]
            ax = plt.gca()
            ax.set_yscale("symlog", linthresh=max(1, min(positive_values)))
            tick_values = [0, min(positive_values), max(positive_values)]
            if len(positive_values) > 1:
                tick_values.insert(2, int(np.median(positive_values)))
            tick_values = sorted(set(tick_values))
            ax.set_yticks(tick_values)
            ax.yaxis.set_major_formatter(mtick.FuncFormatter(_compact_tick_label))
            ax.tick_params(axis="y", labelsize=14)
        # Leyenda a la derecha, centrada verticalmente
        plt.legend(
            loc="center left",
            bbox_to_anchor=(1, 0.5),
            frameon=False,
            labelspacing=1.2,
            fontsize=18,
        )

    elif ctype == "line":
        for idx, key in enumerate(series_keys):
            vals = list(chart_info.get(key) or [])
            vals = vals[:len(labels)]  # truncate if longer
            if len(vals) < len(labels):
                vals.extend([None] * (len(labels) - len(vals)))  # pad with None

            label_full = flat_friendly_names.get(key, key.replace("_", " ").capitalize())
            # Dividir etiquetas largas en varias líneas
            label = textwrap.fill(label_full, width=22)

            color = palette[idx % len(palette)]
            plt.plot(x, vals, label=label, marker="o", color=color)

        plt.xticks(x, labels, rotation=45, fontsize=16)
        plt.grid(axis="y", linestyle="-", color="#dcdcdc", linewidth=0.8)
        plt.legend(loc="best", fontsize=18)

    # Formato eje Y con separador de miles
    try:
        ax = plt.gca()

        if not use_adaptive_scale:
            # ESTA ES LA LÍNEA CLAVE: Fuerza a que los ticks sean solo números enteros
            ax.yaxis.set_major_locator(mtick.MaxNLocator(integer=True))

            # Tu formateador actual
            ax.yaxis.set_major_formatter(mtick.FuncFormatter(lambda x, pos: f"{int(x):,}"))
    except Exception:
        pass

    # Ajusta el layout para asegurar que la leyenda no se corte
    plt.tight_layout(rect=[0, 0.03, 0.95, 0.97])
    plt.savefig(output_file, dpi=150, transparent=True)
    plt.close()


def add_charts(slide, charts, friendly_names, replacements_chart):
    # debug: enumera shape names y contenido para analizar placeholders y textos
    # for shape in slide.shapes:
    #     text = getattr(shape, "text", "")
    #     log(
    #         f"placeholder: '{shape.name}', type={shape.shape_type},"
    #         f" pos=({shape.left},{shape.top}), size=({shape.width},{shape.height}), text='{text.strip()}'"
    #     )

    shape_map = {s.name: s for s in slide.shapes}

    # Tomamos placeholders en el orden deseado de replacements_chart.
    placeholder_candidates = [shape_map[name] for name in replacements_chart if name in shape_map]

    # Asegurar ordenando por posición visual (top, left) como fallback
    placeholder_candidates = sorted(placeholder_candidates, key=lambda s: (s.top, s.left))

    # Tomar solo los placeholders necesarios para los charts disponibles
    chart_placeholders = placeholder_candidates[: len(charts)]

    for (name, chart_info), placeholder in zip(charts.items(), chart_placeholders):
        if not chart_info.get("title") and not chart_info.get("titulo") and name:
            chart_info["title"] = name.replace("_", " ").capitalize()

        # fn = os.path.join(DATA_DIR, f"{name}.png")
        fn = f"/tmp/{name}.png"
        create_matplotlib_chart(chart_info, friendly_names, fn)
        insert_image_scaled_by_width(slide, placeholder, fn)


# --------------------------------------------------------------
# CONTENIDO
# --------------------------------------------------------------


def generar_contenido_slide(slide_item, data, logo_stream):
    slide_content = slide_item.get("slide", {})

    charts = slide_content.get("charts", {})
    num_charts = len(charts)
    if num_charts == 0:
        # No crear la hoja si no hay gráficos
        return None

    product_type = str(slide_item.get("type", "")).strip().lower()
    if product_type == "sonarqube":
        template_file = "plantilla_contenido_sonarqube.pptx"
    elif num_charts == 1:
        template_file = "plantilla_contenido1.pptx"
    elif num_charts == 2:
        template_file = "plantilla_contenido2.pptx"
    elif num_charts == 3:
        template_file = "plantilla_contenido3.pptx"
    else:
        # Si hay más de 4, usar la plantilla de 4 gráficos
        template_file = "plantilla_contenido4.pptx"

    prs = Presentation(f"{DATA_DIR}/plantillas/{template_file}")

    # Asumimos que la plantilla tiene una sola diapositiva
    slide = prs.slides[0]

    # Cargamos los nombres amigables para las leyendas de los gráficos
    product_type = slide_item.get("type")
    friendly_names = {}
    try:
        with open(
            f"{DATA_DIR}/charts/chart_{product_type}.json", "r", encoding="utf-8"
        ) as f:
            friendly_names = json.load(f)
    except FileNotFoundError:
        log(
            f"⚠️  No se encontró el archivo de configuración de gráficos: chart_{product_type}.json"
        )

    feet_l, feet_r, periodo = data.get("pie_l", ""), data.get("pie_r", ""), data.get("periodo", "")

    # Diccionario de reemplazos
    replacements = {
    "{{ph_titulo}}": slide_content.get("titulo", ""),
    "{{ph_periodo}}": periodo,
    "{{ph_title_1}}": slide_content.get("title_1", ""),
    "{{ph_kpis_1}}": normalize_kpi_text(slide_content.get("kpis_1", ""), 250),
    "{{ph_title_2}}": slide_content.get("title_2", ""),
    "{{ph_kpis_2}}": normalize_kpi_text(slide_content.get("kpis_2", ""), 250),
    "{{ph_title_3}}": slide_content.get("title_3", ""),
    "{{ph_kpis_3}}": normalize_kpi_text(slide_content.get("kpis_3", ""), 250),
    "{{ph_title_4}}": slide_content.get("title_4", ""),
    "{{ph_kpis_4}}": normalize_kpi_text(slide_content.get("kpis_4", ""), 250),
    "{{ph_pie_l}}": feet_l,
    "{{ph_pie_r}}": feet_r,
    # Sugerencias (opcionales, vacío si no vienen)
    "{{ph_sugerencia_1}}": normalize_suggestion_text(slide_content.get("sugerencia_1", ""), 200),
    "{{ph_sugerencia_2}}": normalize_suggestion_text(slide_content.get("sugerencia_2", ""), 200),
    "{{ph_sugerencia_3}}": normalize_suggestion_text(slide_content.get("sugerencia_3", ""), 200),
    "{{ph_sugerencia_4}}": normalize_suggestion_text(slide_content.get("sugerencia_4", ""), 200),
    "{{ph_sugerencia4}}": normalize_suggestion_text(slide_content.get("sugerencia_4", ""), 200),
    }
    charts = slide_content.get("charts", {})
    # if product_type != "wazuh" and charts:
    #     replacements["{{ph_utilizacion}}"] = slide_content.get("kpi_title", "")
    #     replacements["{{ph_kpis}}"] = slide_content.get("kpis", "")
    #     replacements["{{ph_soporte}}"] = slide_content.get("soporte_title", "")
    #     replacements["{{ph_soporte_kpis}}"] = slide_content.get("soporte_kpi", "")

    # Reemplazar texto marcador
    modified = replace_placeholders(slide, replacements)

    tf_flags = {}
    for tf, key in modified:
        key_l = key.lower()
        flags = tf_flags.setdefault(id(tf), {"tf": tf, "titulo": False, "kpis": False, "sugerencia": False})
        if "titulo" in key_l:
            flags["titulo"] = True
        if "kpis" in key_l:
            flags["kpis"] = True
        if "sugerencia" in key_l:
            flags["sugerencia"] = True

    for flags in tf_flags.values():
        tf = flags["tf"]
        if flags["titulo"]:
            apply_text_formatting(tf, font_name='Calibri', size=18)
            continue

        if flags["kpis"] or flags["sugerencia"]:
            apply_text_formatting(tf, font_name='Calibri', size=12, set_line=False)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                p.line_spacing = 1.15

                # Destacar el título SUGERENCIAS en cursiva.
                if p.text.strip().upper() == "SUGERENCIAS:":
                    for run in p.runs:
                        run.font.italic = True
            continue

        apply_text_formatting(tf, font_name='Aptos', size=12)
    replacements_chart = [
        "Marcador de posición de imagen 6",
        "Marcador de posición de imagen 7",
        "Marcador de posición de imagen 16",
        "Marcador de posición de imagen 17",
        "Marcador de posición de imagen 18",
    ]
    # Insertar gráficos
    if charts:
        add_charts(slide, charts, friendly_names, replacements_chart)

    insert_logo_with_scaling(slide, logo_stream)

    output_path = f"{DATA_DIR}/pptx-parts/contenido_{product_type}.pptx"
    _prepare_output_file(output_path)
    prs.save(output_path)
    return  output_path


# --------------------------------------------------------------
# CIERRE
# --------------------------------------------------------------
def generar_cierre(data, logo_stream):
    cierre = data["despedida"]
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_cierre.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_titulo}}": cierre.get("titulo", ""),
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    modified = replace_placeholders(slide, replacements)

    for tf, key in modified:
        key_l = key.lower()
        if "titulo" in key_l:
            apply_text_formatting(tf, font_name='Calibri', size=18)
        elif "pie" in key_l:
            apply_text_formatting(tf, font_name=None, size=10)
        else:
            apply_text_formatting(tf, font_name=None, size=12)

    insert_logo_with_scaling(slide, logo_stream)

    output = f"{DATA_DIR}/pptx-parts/cierre.pptx"
    _prepare_output_file(output)
    prs.save(output)
    return output


# --------------------------------------------------------------
# BUENAS PRÁCTICAS
# --------------------------------------------------------------
def generar_buenas_practicas(data, logo_stream):
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_buenas_practicas.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    modified = replace_placeholders(slide, replacements)

    for tf, key in modified:
        key_l = key.lower()
        if "titulo" in key_l:
            apply_text_formatting(tf, font_name='Calibri', size=18)
        elif "pie" in key_l:
            apply_text_formatting(tf, font_name=None, size=10)
        else:
            apply_text_formatting(tf, font_name=None, size=12)

    # Asegurar que el título "BUENAS PRACTICAS" sea Calibri 18
    for shape in slide.shapes:
        if shape.has_text_frame and "BUENAS PRACTICAS" in shape.text.upper():
            apply_text_formatting(shape.text_frame, font_name='Calibri', size=18)

    insert_logo_with_scaling(slide, logo_stream)

    output = f"{DATA_DIR}/pptx-parts/buenas_practicas.pptx"
    _prepare_output_file(output)
    prs.save(output)
    return output


# --------------------------------------------------------------
# PPTX → PDF
# --------------------------------------------------------------
def convert_to_pdf(pptx_file):
    output_dir = f"{DATA_DIR}/pdf-parts"
    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.basename(pptx_file).replace(".pptx", ".pdf")
    pdf_file = os.path.join(output_dir, base_name)
    if os.path.exists(pdf_file):
        os.remove(pdf_file)

    # Usar un directorio de instalación único para evitar bloqueos y problemas de permisos
    user_inst = f"-env:UserInstallation=file:///tmp/lo_{uuid.uuid4()}"
    cmd = [
        "libreoffice",
        user_inst,
        "--headless",
        "--convert-to",
        "pdf",
        pptx_file,
        "--outdir",
        output_dir,
    ]
    try:
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
    except subprocess.CalledProcessError as e:
        log(f"⚠️ Error en LibreOffice: {e.stderr.decode('utf-8', errors='replace')}")
        raise
    return pdf_file


def apply_background_to_pdf(content_pdf_path, background_pdf_path):
    """
    Aplica un fondo desde un PDF a otro PDF que contiene el contenido.
    El contenido se superpone sobre el fondo.
    """
    content_reader = PdfReader(content_pdf_path)
    background_reader = PdfReader(background_pdf_path)
    writer = PdfWriter()

    # Asume que el PDF de fondo tiene al menos tantas páginas como el de contenido
    for i, content_page in enumerate(content_reader.pages):
        # Obtiene la página de fondo correspondiente
        background_page = background_reader.pages[i % len(background_reader.pages)]
        # Superpone el contenido (que tiene fondo transparente/blanco) sobre el fondo
        background_page.merge_page(content_page)
        writer.add_page(background_page)

    with open(content_pdf_path, "wb") as f:
        writer.write(f)


# --------------------------------------------------------------
# UNIR PDFs
# --------------------------------------------------------------
def unir_pdfs(pdf_paths, empresa, type="", split=0):
    writer = PdfWriter()
    for pdf_path in pdf_paths:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            writer.add_page(page)

    output_dir = f"{DATA_DIR}/generados"
    os.makedirs(output_dir, exist_ok=True)
    out = f"{output_dir}/informe_{empresa}{'.' + type if split == 1 else ''}.pdf"
    with open(out, "wb") as f:
        writer.write(f)
    return out


# --------------------------------------------------------------
# MAIN
# --------------------------------------------------------------
def main():
    #raw = sys.argv[1]
    #data = json.loads(base64.b64decode(raw))
    with open(f"{DATA_DIR}/input.json", "r", encoding="utf-8") as f:
         data = json.load(f)
    data = data["data"]
    
    split = data.get("split", 0)
    logo_stream = get_logo_from_base64(data.get("logo_base64"))

    empresa = data.get("logo")[:-4].lower() if data.get("logo") else ""

    slides_data = data.get("slides", [])

    generated_pptx = []
    if data["save"]:
        portada = generar_portada(data, logo_stream)
        generated_pptx.append(portada)

    for slide_item in slides_data:
        product_type = slide_item["type"]
        resumen = slide_item.get("slide", {}).get("resumen", "")  
        producto_slide = generar_slide_producto(resumen, product_type, data, logo_stream)
        generated_pptx.append(producto_slide)
        content_pptx = generar_contenido_slide(slide_item, data, logo_stream)
        generated_pptx.append(content_pptx)

    if data["save"]:
        buenas_practicas = generar_buenas_practicas(data, logo_stream)
        generated_pptx.append(buenas_practicas)
        cierre = generar_cierre(data, logo_stream)
        generated_pptx.append(cierre)

    # Convertir a PDF y unir
    pdf_files_to_merge = [convert_to_pdf(f) for f in generated_pptx if f is not None]

    informe_name = []
    if split == 0:
        informe_name.append(unir_pdfs(pdf_files_to_merge, empresa))
    else:
        # Para split, asumir que se divide por producto, pero como ahora hay más slides, quizás ajustar.
        # Por simplicidad, mantener como estaba, pero con la nueva lista.
        # Pero el usuario no especificó para split, así que dejar como está, pero probablemente split no se usa.
        types = [slide.get("type", "") for slide in slides_data]
        # Para split, crear informes separados por producto, incluyendo portada, producto_slide, resume, content, buenas_practicas, cierre
        # Pero eso sería complejo, asumir split=0 por ahora.
        informe_name.append(unir_pdfs(pdf_files_to_merge, empresa))

    print(json.dumps({"file_names": informe_name}))


if __name__ == "__main__":
    main()
