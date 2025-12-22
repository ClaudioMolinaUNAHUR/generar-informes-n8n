#!/usr/bin/env python3
import sys
import json
import base64
import subprocess
import os
import warnings
from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib.pyplot as plt
from pypdf import PdfReader, PdfWriter
import requests
import textwrap
import numpy as np
from PIL import Image
from io import BytesIO

warnings.filterwarnings("ignore")
DATA_DIR = "/data"


# --------------------------------------------------------------
# UTILS
# --------------------------------------------------------------
def get_logo_from_base64(base64_string: str) -> BytesIO | None:
    """
    Decodifica una cadena Base64 y devuelve un objeto BytesIO con los datos de la imagen.
    Retorna None si la cadena está vacía o es inválida.
    """
    if not base64_string:
        return None
    try:
        image_data = base64.b64decode(base64_string)
        return BytesIO(image_data)
    except base64.binascii.Error:
        log("⚠️ Error de decodificación Base64. La cadena del logo podría ser inválida.")
        return None


def replace_placeholders(slide, replacements):
    """
    Busca un placeholder por su nombre (definido en el Panel de Selección de PowerPoint)
    y reemplaza su texto por el valor correspondiente.
    """

    for key, value in replacements.items():
        for shape in slide.shapes:
            if shape.has_text_frame:
                texto = shape.text.strip()
                if texto == key:
                    # Asegurar que el valor es string
                    if not isinstance(value, str):
                        val_str = str(value)
                    else:
                        val_str = value

                    # Reemplazar secuencias literales "\\n" por saltos de línea reales
                    # y también manejar escapes dobles si vienen.
                    val_str = val_str.replace("\\n", "\n")
                    val_str = val_str.replace("\\\n", "\n")

                    # Asignar al cuadro de texto
                    shape.text = val_str


def insert_image_scaled_by_width(slide, placeholder, image_path_or_stream):
    """
    Reemplaza un placeholder con una imagen, escalándola para que ocupe todo el ancho
    del placeholder y ajustando el alto proporcionalmente. La imagen se centra verticalmente.
    """
    # 1. Obtener dimensiones y posición del placeholder
    ph_left, ph_top = placeholder.left, placeholder.top
    ph_width, ph_height = placeholder.width, placeholder.height

    # 2. Eliminar el placeholder original para evitar duplicados
    placeholder.element.getparent().remove(placeholder.element)

    # 3. Insertar la imagen con el ancho del placeholder.
    #    python-pptx ajustará automáticamente el alto para mantener la proporción.
    pic = slide.shapes.add_picture(
        image_path_or_stream, ph_left, ph_top, width=ph_width
    )

    # 4. Centrar la imagen verticalmente en el espacio del placeholder original
    new_height = pic.height
    pic.top = ph_top + (ph_height - new_height) // 2


def insert_logo_preserving_aspect(slide, placeholder, logo_stream):
    # área disponible
    ph_left, ph_top = placeholder.left, placeholder.top
    ph_w, ph_h = placeholder.width, placeholder.height

    # borrar placeholder original
    placeholder.element.getparent().remove(placeholder.element)

    # insertar imagen sin escalar
    pic = slide.shapes.add_picture(logo_stream, ph_left, ph_top)

    # tamaño real de la imagen
    img_w, img_h = pic.width, pic.height

    # calcular factor manteniendo aspecto
    scale = min(ph_w / img_w, ph_h / img_h)

    # aplicar tamaño escalado
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)

    pic.width = new_w
    pic.height = new_h

    # centrar dentro del placeholder original
    pic.left = ph_left + (ph_w - new_w) // 2
    pic.top = ph_top + (ph_h - new_h) // 2


def _insert_logo_with_scaling(slide, logo_stream):
    """
    Busca el primer placeholder de tipo 'Picture' (18) e inserta el logo
    dentro de sus límites, manteniendo la relación de aspecto y eliminando el placeholder original.
    """
    # 18 es el tipo de placeholder para 'Picture'
    LOGO_PLACEHOLDER_TYPE = 18

    if not logo_stream:
        return

    for shape in slide.placeholders:
        if shape.placeholder_format.type == LOGO_PLACEHOLDER_TYPE:

            insert_logo_preserving_aspect(slide, shape, logo_stream)
            break  # Solo insertamos el primer logo que encontremos


def log(msg):
    print(msg, flush=True)


def set_placeholder_text(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            if ph.has_text_frame:
                ph.text = text
            return
    # Si no existe, opcionalmente crear un cuadro de texto
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = tb.text_frame
    tf.text = text

def create_composite_logo_from_base64_list(
    logos_base64_list: list[str],
    target_height: int = 120,  # altura uniforme
) -> BytesIO | None:

    if not logos_base64_list:
        return None

    images = []
    for b64_string in logos_base64_list:
        img_stream = get_logo_from_base64(b64_string)
        if img_stream:
            try:
                img = Image.open(img_stream).convert("RGBA")

                # 🔥 resize proporcional
                ratio = target_height / img.height
                new_width = int(img.width * ratio)
                img = img.resize((new_width, target_height), Image.LANCZOS)

                images.append(img)
            except Exception as e:
                log(f"⚠️ Error opening image from stream: {e}")

    if not images:
        return None

    max_width = max(img.width for img in images)
    total_height = sum(img.height for img in images)

    composite_image = Image.new("RGBA", (max_width, total_height), (0, 0, 0, 0))

    y_offset = 0
    for img in images:
        x_offset = (max_width - img.width) // 2
        composite_image.paste(img, (x_offset, y_offset), img)
        y_offset += img.height

    output_stream = BytesIO()
    composite_image.save(output_stream, format="PNG")
    output_stream.seek(0)

    return output_stream



# --------------------------------------------------------------
# PORTADA
# --------------------------------------------------------------
def generar_portada(data, logo_stream):
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_portada.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_titulo}}": data.get("titulo_portada", ""),
        "{{ph_subtitle}}": data.get("subtitulo_portada", ""),
        "{{ph_fecha}}": data.get("fecha_portada", ""),
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    replace_placeholders(slide, replacements)

    # Busca un placeholder de tipo imagen (18) para el logo.
    _insert_logo_with_scaling(slide, logo_stream)

    output = f"{DATA_DIR}/pptx-parts/portada.pptx"
    prs.save(output)
    return output


# --------------------------------------------------------------
# CIERRE
# --------------------------------------------------------------
def generar_cierre(data, logo_stream):
    cierre = data["despedida"]
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_cierre.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_titulo}}": cierre.get("titulo", ""),
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    replace_placeholders(slide, replacements)

    _insert_logo_with_scaling(slide, logo_stream)

    output = f"{DATA_DIR}/pptx-parts/cierre.pptx"
    prs.save(output)
    return output


# --------------------------------------------------------------
# PPTX → PDF
# --------------------------------------------------------------
def convert_to_pdf(pptx_file):
    output_dir = f"{DATA_DIR}/pdf-parts"
    base_name = os.path.basename(pptx_file).replace(".pptx", ".pdf")
    pdf_file = os.path.join(output_dir, base_name)
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        pptx_file,
        "--outdir",
        output_dir,
    ]
    subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
    return pdf_file


def apply_background_to_pdf(content_pdf_path, background_pdf_path):
    """
    Aplica un fondo desde un PDF a otro PDF que contiene el contenido.
    El contenido se superpone sobre el fondo.
    """
    content_reader = PdfReader(content_pdf_path)
    background_reader = PdfReader(background_pdf_path)
    writer = PdfWriter()

    # Asume que el PDF de fondo tiene al menos tantas páginas como el de contenido
    for i, content_page in enumerate(content_reader.pages):
        # Obtiene la página de fondo correspondiente
        background_page = background_reader.pages[i % len(background_reader.pages)]
        # Superpone el contenido (que tiene fondo transparente/blanco) sobre el fondo
        background_page.merge_page(content_page)
        writer.add_page(background_page)

    with open(content_pdf_path, "wb") as f:
        writer.write(f)


# --------------------------------------------------------------
# UNIR PDFs
# --------------------------------------------------------------
def unir_pdfs(pdf_paths, empresa):
    writer = PdfWriter()
    for pdf_path in pdf_paths:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            writer.add_page(page)
    out = f"{DATA_DIR}/generados/informe_{empresa}.pdf"
    with open(out, "wb") as f:
        writer.write(f)
    return out


# --------------------------------------------------------------
# MAIN
# --------------------------------------------------------------
def main():
    raw = sys.argv[1]
    input_data = json.loads(base64.b64decode(raw))

    main_data = input_data["main"]
    emp_codes = input_data.get("emp_codes", [])
    logos_base64_list = input_data.get("logos_base64", [])

    logo_stream = create_composite_logo_from_base64_list(logos_base64_list)

    empresa = ""
    length_emp_codes = len(emp_codes)
    for i, emp_code in enumerate(emp_codes):
        empresa += (
            emp_code + "-"
            if i != length_emp_codes - 1
            else emp_code
        )
    empresa = empresa.lower()
    portada_pptx_file = generar_portada(main_data, logo_stream)
    cierre_pptx_file = generar_cierre(main_data, logo_stream)

    pdf_files_to_merge = []
    pdf_files_to_merge.append(convert_to_pdf(portada_pptx_file))

    full_informes_paths = [
        os.path.join(DATA_DIR, "generados", f"informe_{f.lower()}.pdf") for f in emp_codes
    ]
    pdf_files_to_merge.extend(full_informes_paths)

    pdf_files_to_merge.append(convert_to_pdf(cierre_pptx_file))

    final_pdf = unir_pdfs(pdf_files_to_merge, empresa)

    # with open(final_pdf, "rb") as f:
    #     b64 = base64.b64encode(f.read()).decode()

    print(json.dumps({"file_name": os.path.basename(f"informe_{empresa}")}))


if __name__ == "__main__":
    main()
