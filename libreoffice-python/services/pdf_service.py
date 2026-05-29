import os
import json
import uuid
import subprocess
import shutil
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from utils.helpers import (
    DATA_DIR,
    log,
    insert_logo_with_scaling,
    apply_text_formatting,
    normalize_kpi_text,
    normalize_suggestion_text,
    insert_image_scaled_by_width,
)
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Pt
import re
from services.chart_service import create_matplotlib_chart, add_charts
from pptx.dml.color import RGBColor

def _prepare_output_file(path):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    if os.path.exists(path):
        os.remove(path)




def replace_placeholders(slide, replacements):
    """
    Busca placeholders en los text_frames de la slide y los reemplaza.

    Mejoras:
    - Reemplazo robusto a nivel XML: reconstruye el texto completo del párrafo
      concatenando todos sus runs antes de buscar placeholders, evitando fallos
      cuando un placeholder está partido entre varios runs.
    - Negrita automática: en bloques de texto plano, cualquier fragmento que
      precede a ':' dentro de una línea se renderiza en negrita (ej. "Módulo: texto").
    - Limpieza de placeholders no resueltos: elimina párrafos completos cuyo
      texto sea únicamente un placeholder sin valor asignado, en lugar de dejar
      la línea vacía o con el token visible.
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    modified = []
    unresolved_pattern = re.compile(r"\{\{ph_[^}]+\}\}")

    def _value_to_text(value):
        if isinstance(value, list):
            return "".join(
                str(item.get("text", "")) if isinstance(item, dict) else str(item)
                for item in value
            )
        return str(value) if value is not None else ""

    def _rebuild_para_text(para):
        """Devuelve el texto completo del párrafo uniendo todos sus runs."""
        return "".join(run.text or "" for run in para.runs)

    def _set_para_with_bold_labels(para, text, base_font_size=None, base_font_name=None):
        """
        Reemplaza el contenido de un párrafo con soporte de negrita automática.
        Si la línea tiene el patrón 'Etiqueta: resto', pone 'Etiqueta:' en negrita.
        Respeta saltos de línea (\n) creando nuevos párrafos en el text_frame padre.
        Retorna lista de párrafos adicionales creados (para alineación posterior).
        """
        p_elem = para._p
        tf = para._p.getparent()

        # Limpiar runs existentes del párrafo
        for r in p_elem.findall(qn("a:r")):
            p_elem.remove(r)

        lines = text.split("\n")
        extra_paras = []

        def _add_run_to_p(p_el, run_text, bold=False, font_size=None, font_name=None):
            """Agrega un run XML al elemento párrafo."""
            r_el = etree.SubElement(p_el, qn("a:r"))
            rPr = etree.SubElement(r_el, qn("a:rPr"), attrib={"lang": "es-AR", "dirty": "0"})
            if bold:
                rPr.set("b", "1")
            if font_size:
                rPr.set("sz", str(int(font_size * 100)))
            if font_name:
                latin = etree.SubElement(rPr, qn("a:latin"))
                latin.set("typeface", font_name)
            t_el = etree.SubElement(r_el, qn("a:t"))
            t_el.text = run_text

        def _write_line_to_p(p_el, line):
            """Escribe una línea con negrita automática antes del primer ':'."""
            colon_idx = line.find(":")
            if colon_idx > 0:
                label = line[: colon_idx + 1]   # incluye el ':'
                rest  = line[colon_idx + 1 :]   # todo lo que sigue
                _add_run_to_p(p_el, label, bold=True,  font_size=base_font_size, font_name=base_font_name)
                if rest:
                    _add_run_to_p(p_el, rest, bold=False, font_size=base_font_size, font_name=base_font_name)
            else:
                _add_run_to_p(p_el, line, bold=False, font_size=base_font_size, font_name=base_font_name)

        # Primera línea va al párrafo original
        _write_line_to_p(p_elem, lines[0])

        # Líneas adicionales: insertar nuevos párrafos después del original
        insert_after = p_elem
        for line in lines[1:]:
            new_p = copy.deepcopy(p_elem)
            # Limpiar runs del nuevo párrafo clonado
            for r in new_p.findall(qn("a:r")):
                new_p.remove(r)
            _write_line_to_p(new_p, line)
            insert_after.addnext(new_p)
            insert_after = new_p
            extra_paras.append(new_p)

        return extra_paras

    # ── 1) Reemplazo robusto: reconstruye texto por párrafo y reemplaza ──────
    already_processed = set()  # shape ids ya procesados en modo exacto

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame

        # Log de depuración para ver qué estamos procesando
        log(f"DEBUG: replace_placeholders - Procesando shape: '{shape.name}', Texto: '{tf.text.strip()}'")

        # Texto completo del cuadro para detectar si es reemplazo "exacto"
        full_box_text = tf.text.strip()

        # Modo exacto: el cuadro contiene SOLO un placeholder
        if full_box_text in replacements:
            key = full_box_text
            value = replacements[key]
            already_processed.add(id(shape))

            # Limpiar todos los párrafos salvo el primero
            while len(tf.paragraphs) > 1:
                p_to_remove = tf.paragraphs[-1]._p
                p_to_remove.getparent().remove(p_to_remove)

            first_para = tf.paragraphs[0]

            if isinstance(value, list):
                # Lista de dicts con {text, bold}
                for r in first_para._p.findall(qn("a:r")):
                    first_para._p.remove(r)
                insert_after = first_para._p
                for item in value:
                    if not isinstance(item, dict):
                        continue
                    new_p = copy.deepcopy(first_para._p)
                    for r in new_p.findall(qn("a:r")):
                        new_p.remove(r)
                    r_el = etree.SubElement(new_p, qn("a:r"))
                    rPr = etree.SubElement(r_el, qn("a:rPr"), attrib={"lang": "es-AR", "dirty": "0"})
                    if item.get("bold", False):
                        rPr.set("b", "1")
                    t_el = etree.SubElement(r_el, qn("a:t"))
                    t_el.text = item.get("text", "")
                    insert_after.addnext(new_p)
                    insert_after = new_p
            else:
                val_str = _value_to_text(value).replace("\\n", "\n").replace("\\\n", "\n")
                _set_para_with_bold_labels(first_para, val_str)

            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.JUSTIFY
            modified.append((tf, key))
            continue

        # Modo embebido: el cuadro tiene texto con uno o más placeholders intercalados
        changed = False
        replaced_keys_here = []

        for para in tf.paragraphs:
            para_text = _rebuild_para_text(para)
            new_para_text = para_text

            for key, value in replacements.items():
                if key in new_para_text:
                    val_str = _value_to_text(value).replace("\\n", "\n").replace("\\\n", "\n")
                    new_para_text = new_para_text.replace(key, val_str)
                    if key not in replaced_keys_here:
                        replaced_keys_here.append(key)

            if new_para_text != para_text:
                _set_para_with_bold_labels(para, new_para_text)
                para.alignment = PP_ALIGN.JUSTIFY
                changed = True

        if changed:
            # Log resumido del cambio
            log(f"DEBUG: replace_placeholders - Modificada forma '{shape.name}' (reemplazo embebido)")
            for key in replaced_keys_here:
                modified.append((tf, key))

    # ── 2a) Limpieza: párrafos vacíos tras reemplazo con "" ──────────────────
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        paras_to_remove = []
        for para in tf.paragraphs:
            para_text = _rebuild_para_text(para).strip()
            if not para_text:
                paras_to_remove.append(para._p)
        for p_elem in paras_to_remove:
            log(f"DEBUG: replace_placeholders - Limpiando párrafo vacío en shape '{shape.name}'")
            parent = p_elem.getparent()
            if parent is not None:
                siblings = parent.findall(qn("a:p"))
                if len(siblings) > 1:
                    parent.remove(p_elem)

    # ── 2) Limpieza: eliminar párrafos con placeholders no resueltos ──────────
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        paras_to_remove = []

        for para in tf.paragraphs:
            para_text = _rebuild_para_text(para)
            if unresolved_pattern.search(para_text):
                cleaned = unresolved_pattern.sub("", para_text).strip()
                if not cleaned:
                    paras_to_remove.append(para._p)

        for p_elem in paras_to_remove:
            parent = p_elem.getparent()
            if parent is not None:
                siblings = parent.findall(qn("a:p"))
                if len(siblings) > 1:
                    parent.remove(p_elem)
                else:
                    # Limpiar runs del único párrafo
                    for r in p_elem.findall(qn("a:r")):
                        p_elem.remove(r)

    return modified



def convert_to_pdf(pptx_file: str) -> str:
    output_dir = f"{DATA_DIR}/pdf-parts"
    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.basename(pptx_file).replace(".pptx", ".pdf")
    pdf_file = os.path.join(output_dir, base_name)
    if os.path.exists(pdf_file):
        os.remove(pdf_file)

    profile_dir = f"/tmp/lo_{uuid.uuid4()}"
    cmd = [
        "libreoffice",
        f"-env:UserInstallation=file://{profile_dir}",
        "--headless",
        "--convert-to",
        "pdf",
        pptx_file,
        "--outdir",
        output_dir,
    ]

    try:
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
    except subprocess.CalledProcessError as e:
        log(f"⚠️ Error en LibreOffice: {e.stderr.decode('utf-8', errors='replace')}")
        raise Exception(f"LibreOffice error: {e.stderr.decode()}")
    finally:
        if os.path.exists(profile_dir):
            shutil.rmtree(profile_dir, ignore_errors=True)

    return pdf_file


def generar_portada(data, logo_stream):
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_portada.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_titulo}}": data.get("titulo_portada", ""),
        "{{ph_subtitle}}": data.get("subtitulo_portada", ""),
        "{{ph_fecha}}": data.get("fecha_portada", ""),
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    # Insertar logo buscando el marcador {{ph_logo}}
    insert_logo_with_scaling(slide, logo_stream)

    # Ejecutamos el reemplazo normal
    modified = replace_placeholders(slide, replacements)

    # Aplicar formato SOLO a los campos necesarios
    for tf, key in modified:
        key_l = key.lower()
        if "titulo" in key_l or "subtitle" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=18)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.line_spacing = 1.5
        elif "fecha" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=12)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.line_spacing = 1.2
        elif "pie" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=10)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT if "l" in key_l else PP_ALIGN.RIGHT
                p.line_spacing = 1.0

    output = f"{DATA_DIR}/pptx-parts/portada.pptx"
    _prepare_output_file(output)
    prs.save(output)
    return output


def generar_contenido_slide(slide_item, data, logo_stream):
    slide_content = slide_item.get("slide", {})

    charts = slide_content.get("charts", {})
    num_charts = len(charts)
    if num_charts == 0:
        # No crear la hoja si no hay gráficos
        return None
    elif num_charts == 1:
        template_file = "plantilla_contenido_1.pptx"
    elif num_charts == 2:
        template_file = "plantilla_contenido_2.pptx"
    elif num_charts == 3:
        template_file = "plantilla_contenido_3.pptx"
    else:
        # Si hay más de 4, usar la plantilla de 4 gráficos
        template_file = "plantilla_contenido_4.pptx"
    print(f"Generando slide de contenido para {slide_item.get('type')} con {num_charts} gráficos usando plantilla {template_file}")

    prs = Presentation(f"{DATA_DIR}/plantillas/{template_file}")

    # Asumimos que la plantilla tiene una sola diapositiva
    slide = prs.slides[0]

    # Cargamos los nombres amigables para las leyendas de los gráficos
    product_type = slide_item.get("type")
    friendly_names = {}
    try:
        with open(
            f"{DATA_DIR}/charts/chart_{product_type}.json", "r", encoding="utf-8"
        ) as f:
            friendly_names = json.load(f)
    except FileNotFoundError:
        log(
            f"⚠️  No se encontró el archivo de configuración de gráficos: chart_{product_type}.json"
        )

    feet_l, feet_r, periodo = (
        data.get("pie_l", ""),
        data.get("pie_r", ""),
        data.get("periodo", ""),
    )

    # Diccionario de reemplazos
    desc_val = str(slide_content.get("desc", "") or "")
    replacements = {
        "{{ph_titulo}}": slide_content.get("titulo", ""),
        "{{ph_periodo}}": periodo,
        "{{ph_title_1}}": slide_content.get("title_1", ""),
        "{{ph_kpis_1}}": normalize_kpi_text(slide_content.get("kpis_1", ""), 250),
        "{{ph_title_2}}": slide_content.get("title_2", ""),
        "{{ph_kpis_2}}": normalize_kpi_text(slide_content.get("kpis_2", ""), 250),
        "{{ph_title_3}}": slide_content.get("title_3", ""),
        "{{ph_kpis_3}}": normalize_kpi_text(slide_content.get("kpis_3", ""), 250),
        "{{ph_title_4}}": slide_content.get("title_4", ""),
        "{{ph_kpis_4}}": normalize_kpi_text(slide_content.get("kpis_4", ""), 250),
        "{{ph_pie_l}}": feet_l,
        "{{ph_pie_r}}": feet_r,
        # Asegura que desc nunca sea None
        "{{ph_desc}}": desc_val,
        # Sugerencias (opcionales, vacío si no vienen)
        "{{ph_sugerencia_1}}": normalize_suggestion_text(
            slide_content.get("sugerencia_1", ""), 200
        ),
        "{{ph_sugerencia_2}}": normalize_suggestion_text(
            slide_content.get("sugerencia_2", ""), 200
        ),
        "{{ph_sugerencia_3}}": normalize_suggestion_text(
            slide_content.get("sugerencia_3", ""), 200
        ),
        "{{ph_sugerencia_4}}": normalize_suggestion_text(
            slide_content.get("sugerencia_4", ""), 200
        ),
    }
    print(f"Reemplazos para slide de contenido ({product_type}): {replacements}")
    print(f"Valor insertado en {{ph_desc}}: {desc_val}")
    
    # Insertar gráficos y logo ANTES del reemplazo de texto:
    # Esto evita que la lógica de limpieza de replace_placeholders elimine los marcadores {{ph_...}}
    if charts:
        add_charts(slide, charts, friendly_names)

    log(
        f"DEBUG: generar_contenido_slide - type={product_type}, logo_stream={'sí' if logo_stream else 'no'}"
    )
    logo_inserted = insert_logo_with_scaling(slide, logo_stream)
    log(f"DEBUG: generar_contenido_slide - logo_inserted={logo_inserted}")

    # Reemplazar texto marcador
    modified = replace_placeholders(slide, replacements)

    tf_flags = {}
    for tf, key in modified:
        key_l = key.lower()
        flags = tf_flags.setdefault(
            id(tf), {"tf": tf, "titulo": False, "kpis": False, "sugerencia": False}
        )
        if "titulo" in key_l:
            flags["titulo"] = True
        if "kpis" in key_l:
            flags["kpis"] = True
        if "sugerencia" in key_l:
            flags["sugerencia"] = True

    for flags in tf_flags.values():
        tf = flags["tf"]
        if flags["titulo"]:
            apply_text_formatting(tf, font_name="Aptos", size=18)
            continue

        if flags["sugerencia"] or flags["kpis"] or "title" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=10, set_line=False)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                p.line_spacing = 1.15

                # Destacar el título SUGERENCIAS en cursiva.
                if p.text.strip().upper() == "SUGERENCIAS:":
                    for run in p.runs:
                        run.font.italic = True
            continue

        # Cambia el tamaño de fuente solo para axur y plantilla_axur.pptx
        font_size = 12
        # Buscar el nombre de la plantilla usada
        template_file = None
        # Buscar en la pila de llamadas la variable template_file
        import inspect
        for frame_info in inspect.stack():
            if 'template_file' in frame_info.frame.f_locals:
                template_file = frame_info.frame.f_locals['template_file']
                break
        if slide_item.get("type", "").lower() == "axur" and template_file == "plantilla_axur.pptx":
            font_size = 10
        apply_text_formatting(tf, font_name="Aptos", size=font_size)

    output_path = f"{DATA_DIR}/pptx-parts/contenido_{product_type}.pptx"
    _prepare_output_file(output_path)
    prs.save(output_path)
    return output_path


def generar_slide_producto(
    resumen, product_type, data, logo_stream, pie_l="", pie_r=""
):
    # Normalizar product_type para selección de plantilla y campos
    # Ejemplo: invgate.asj -> invgate
    base_type = product_type.lower().split('.')[0]
    if base_type == "invgate" or base_type == "asj": base_type = "invgate"

    template_path = f"{DATA_DIR}/plantillas/plantilla_{base_type}.pptx"
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Definir campos por tipo de producto
    campos_por_tipo = {
        "uas": ["usu_per", "usu_esp", "solicitudes", "revalida"],
        "beyondtrust": ["pra", "rs", "ps", "adb", "epm"],
        "whalemate": ["sim", "aca", "ana", "grh", "cad"],
        "wazuh": ["ddv", "snc", "enc", "ecn", "iav"], # Incluir ambos por compatibilidad
        "invgate": ["isd", "iam"],
        "invgate.asj": ["isd", "iam"],
        "axur": ["pdm", "th", "cti", "fdd", "ddw", "tkd"],
        "akurtech": ["tra", "pgs", "lgn", "rfg", "alr", "bwl", "rdc", "rdt", "rdp", "rdl"],
    }

    # Mapear product_type a grupo
    tipo_grupo = base_type

    replacements = {
        "{{ph_resumen}}": resumen,
        "{{ph_resume}}": resumen, # Compatibilidad con plantillas en inglés/mixtas
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    # Agregar "VERSION:" como línea antes del resumen en el placeholder
    # Buscar versión en distintas claves posibles del data dict
    version_line = (
        data.get("version") or
        data.get("version_produccion") or
        data.get("produccion_version") or
        ""
    )
    if version_line:
        res_text = f"VERSION: {version_line}\n\n{resumen}"
        replacements["{{ph_resumen}}"] = res_text
        replacements["{{ph_resume}}"] = res_text
    
    # Si el resumen ya trae info de versión embebida, agregar encabezado "VERSION:" igual
    elif resumen and any(kw in resumen.lower() for kw in ["versión", "version", "v."]):
        res_text = f"VERSION:\n{resumen}"
        replacements["{{ph_resumen}}"] = res_text
        replacements["{{ph_resume}}"] = res_text

    # Agregar los campos nuevos si corresponden
    if tipo_grupo:
        for campo in campos_por_tipo[tipo_grupo]:
            replacements[f"{{{{ph_{campo}}}}}"] = data.get(campo, "")

    # Quitar negritas heredadas del template ANTES de reemplazar,
    # para no pisar la negrita que _set_para_with_bold_labels va a aplicar.
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = False

    log(
        f"DEBUG: generar_slide_producto - product_type={product_type}, logo_stream={'sí' if logo_stream else 'no'}"
    )
    logo_inserted = insert_logo_with_scaling(slide, logo_stream)
    log(f"DEBUG: generar_slide_producto - logo_inserted={logo_inserted}")

    modified = replace_placeholders(slide, replacements)

    # Campos de módulos llevan negrita automática en la etiqueta (antes del ':').
    # apply_text_formatting NO debe tocar bold en estos campos.
    CAMPOS_CON_NEGRITA = {
        "usu_per", "usu_esp", "solicitudes", "revalida",
        "pra", "rs", "ps", "adb", "epm",
        "sim", "aca", "ana", "grh", "cad",
        "ddv", "snc", "enc", "iav",
        "isd", "iam", "resumen",
    }

    for tf, key in modified:
        key_l = key.lower()
        campo = key_l.replace("{{ph_", "").replace("}}", "")

        if "titulo" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=18)
        elif "sub" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=14)
        elif "kpis" in key_l or "sugerencia" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=12, set_line=False)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(0)
                p.space_after = Pt(0)
        elif "pie" in key_l:
            apply_text_formatting(tf, font_name=None, size=10)
        elif campo in CAMPOS_CON_NEGRITA:
            # Aplicar fuente y tamaño manualmente SIN tocar bold,
            # para preservar la negrita de las etiquetas (antes del ':').
            # Para axur y akurtech se usa tamaño 11 en lugar de 12.
            from pptx.util import Pt as _Pt
            campo_font_size = 11 if base_type in ("axur", "akurtech") else 12
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.JUSTIFY
                p.line_spacing = 1.5
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = _Pt(campo_font_size)
                    # run.font.bold intacto → preserva lo de _set_para_with_bold_labels
        else:
            # Para axur y akurtech se usa tamaño 11 en lugar de 12.
            campo_font_size = 11 if base_type in ("axur", "akurtech") else 12
            apply_text_formatting(tf, font_name="Aptos", size=campo_font_size)

    output = f"{DATA_DIR}/pptx-parts/producto_{product_type}.pptx"
    _prepare_output_file(output)
    prs.save(output)
    return output


def generar_cierre(data, logo_stream):
    cierre = data["despedida"]
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_cierre.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_titulo}}": cierre.get("titulo", "").replace("\n", " ").replace("\r", " "),
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    # Insertar logo
    insert_logo_with_scaling(slide, logo_stream)

    modified = replace_placeholders(slide, replacements)

    # Aplicar formato SOLO a los campos necesarios
    for tf, key in modified:
        key_l = key.lower()
        if "titulo" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=18, set_line=True)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.line_spacing = 1.5
                for run in p.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255) 
                # Asegurar título
        elif "pie" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=10)
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT if "l" in key_l else PP_ALIGN.RIGHT
                p.line_spacing = 1.0

    output = f"{DATA_DIR}/pptx-parts/cierre.pptx"
    _prepare_output_file(output)
    prs.save(output)
    return output


def generar_buenas_practicas(data, logo_stream):
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_buenas_practicas.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    # Insertar logo
    insert_logo_with_scaling(slide, logo_stream)

    modified = replace_placeholders(slide, replacements)

    for tf, key in modified:
        key_l = key.lower()
        if "titulo" in key_l:
            apply_text_formatting(tf, font_name="Aptos", size=18)
        elif "pie" in key_l:
            apply_text_formatting(tf, font_name=None, size=10)
        else:
            apply_text_formatting(tf, font_name=None, size=12)

    # Asegurar que el título "BUENAS PRACTICAS" sea Aptos 18
    for shape in slide.shapes:
        if shape.has_text_frame and "BUENAS PRACTICAS" in shape.text.upper():
            apply_text_formatting(shape.text_frame, font_name="Aptos", size=18)

    output = f"{DATA_DIR}/pptx-parts/buenas_practicas.pptx"
    _prepare_output_file(output)
    prs.save(output)
    return output


def unir_pdfs(pdf_paths, empresa, type="", split=0):
    writer = PdfWriter()
    for pdf_path in pdf_paths:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            writer.add_page(page)

    output_dir = f"{DATA_DIR}/generados"
    os.makedirs(output_dir, exist_ok=True)
    out = f"{output_dir}/informe_{empresa}{'.' + type if split == 1 else ''}.pdf"
    with open(out, "wb") as f:
        writer.write(f)
    return out