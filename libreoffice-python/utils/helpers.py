import os
import base64
from io import BytesIO
import textwrap
from PIL import Image
from pptx.util import Inches, Pt
import warnings

warnings.filterwarnings("ignore")

DATA_DIR = "/data"

def get_safe_path(relative_path: str) -> str:
    """
    Construye una ruta absoluta y valida que se encuentre dentro de DATA_DIR
    para prevenir ataques de acceso a archivos no autorizados.
    """
    clean_rel_path = relative_path.lstrip("/")
    full_path = os.path.abspath(os.path.join(DATA_DIR, clean_rel_path))
    if not full_path.startswith(os.path.abspath(DATA_DIR)):
        raise ValueError(
            "Acceso denegado: Intento de acceso fuera del directorio de datos."
        )
    return full_path


def read_data_file(relative_path: str) -> bytes:
    path = get_safe_path(relative_path)
    with open(path, "rb") as f:
        return f.read()


def save_data_file(relative_path: str, content: bytes):
    path = get_safe_path(relative_path)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as f:
        f.write(content)


def truncate_text(text, max_chars=300):
    if len(text) > max_chars:
        return text[:max_chars].rsplit(" ", 1)[0] + "…"
    return text


def normalize_kpi_text(text, max_chars=250):
    val = truncate_text((text or "").strip(), max_chars)
    return val if val else " "


def normalize_suggestion_text(text, max_chars=200, line_width=46):
    raw = truncate_text((text or "").strip(), max_chars)
    if not raw:
        return " "

    # Evita bloques visualmente amontonados en una sola línea.
    wrapped_lines = []
    for line in raw.replace("\\n", "\n").splitlines():
        line = line.strip()
        if not line:
            continue
        bullet = line.startswith("•") or line.startswith("-")
        content = line.lstrip("•- ").strip() if bullet else line
        if not content:
            continue
        if bullet:
            wrapped = textwrap.fill(
                content, width=line_width, initial_indent="• ", subsequent_indent="  "
            )
        else:
            wrapped = textwrap.fill(content, width=line_width)
        wrapped_lines.append(wrapped)

    body = "\n".join(wrapped_lines) if wrapped_lines else " "
    if body.strip() and not body.upper().startswith("SUGERENCIAS"):
        return f"SUGERENCIAS:\n{body}"
    return body


def apply_text_formatting(text_frame, font_name="Aptos", size=None, set_line=True):
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        if set_line:
            paragraph.line_spacing = 1.5
        for run in paragraph.runs:
            if font_name is not None:
                run.font.name = font_name
            if size is not None:
                run.font.size = Pt(size)


def log(msg: str):
    print(msg, flush=True)


def get_logo_from_base64(base64_string: str) -> BytesIO | None:
    """
    Decodifica una cadena Base64 y devuelve un objeto BytesIO con los datos de la imagen.
    Retorna None si la cadena está vacía o es inválida.
    """
    if not base64_string:
        return None
    try:
        image_data = base64.b64decode(base64_string)
        return BytesIO(image_data)
    except base64.binascii.Error:
        log(
            "⚠️ Error de decodificación Base64. La cadena del logo podría ser inválida."
        )
        return None


def set_placeholder_text(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            if ph.has_text_frame:
                ph.text = text
                apply_text_formatting(ph.text_frame, font_name="Aptos", size=11)
            return
    # Si no existe, opcionalmente crear un cuadro de texto
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = tb.text_frame
    tf.text = text
    apply_text_formatting(tf, font_name="Aptos", size=11)


def create_composite_logo_from_base64_list(
    logos_base64_list: list[str], target_height: int = 120
) -> BytesIO | None:
    if not logos_base64_list:
        return None

    images = []
    for b64_string in logos_base64_list:
        img_stream = get_logo_from_base64(b64_string)
        if img_stream:
            try:
                img = Image.open(img_stream).convert("RGBA")
                ratio = target_height / img.height
                new_width = int(img.width * ratio)
                img = img.resize((new_width, target_height), Image.LANCZOS)
                images.append(img)
            except Exception as e:
                log(f"⚠️ Error opening image from stream: {e}")

    if not images:
        return None

    max_width = max(img.width for img in images)
    total_height = sum(img.height for img in images)
    composite_image = Image.new("RGBA", (max_width, total_height), (0, 0, 0, 0))

    y_offset = 0
    for img in images:
        x_offset = (max_width - img.width) // 2
        composite_image.paste(img, (x_offset, y_offset), img)
        y_offset += img.height

    output_stream = BytesIO()
    composite_image.save(output_stream, format="PNG")
    output_stream.seek(0)
    return output_stream


def insert_image_scaled_by_width(slide, placeholder, image_path_or_stream):
    """
    Reemplaza un placeholder con una imagen, escalándola para que ocupe todo el ancho
    del placeholder y ajustando el alto proporcionalmente. La imagen se centra verticalmente.
    """
    parent = placeholder.element.getparent()
    if parent is None:
        return  # ya fue eliminado, evitar crash

    ph_left, ph_top = placeholder.left, placeholder.top
    ph_width, ph_height = placeholder.width, placeholder.height

    parent.remove(placeholder.element)

    log(
        f"DEBUG: insert_image_scaled_by_width - placeholder '{placeholder.name}' dims=({ph_width},{ph_height}) pos=({ph_left},{ph_top})"
    )
    pic = slide.shapes.add_picture(image_path_or_stream, ph_left, ph_top, width=ph_width)
    log(
        f"DEBUG: insert_image_scaled_by_width - picture inserted dims=({pic.width},{pic.height}) pos=({pic.left},{pic.top})"
    )

    # Ajustar la imagen para que encaje dentro del placeholder sin desbordar.
    if pic.height > ph_height:
        scale = ph_height / pic.height
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        log(
            f"DEBUG: insert_image_scaled_by_width - picture redimensionada a dims=({pic.width},{pic.height}) para caber en el placeholder"
        )

    pic.left = ph_left + max(0, (ph_width - pic.width) // 2)
    pic.top = ph_top + max(0, (ph_height - pic.height) // 2)
    log(
        f"DEBUG: insert_image_scaled_by_width - picture centered pos=({pic.left},{pic.top})"
    )


def insert_logo_preserving_aspect(slide, placeholder, logo_stream):
    # área disponible
    ph_left, ph_top = placeholder.left, placeholder.top
    ph_w, ph_h = placeholder.width, placeholder.height

    # borrar placeholder original
    placeholder.element.getparent().remove(placeholder.element)

    # insertar imagen sin escalar
    pic = slide.shapes.add_picture(logo_stream, ph_left, ph_top)

    # tamaño real de la imagen
    img_w, img_h = pic.width, pic.height

    # calcular factor manteniendo aspecto
    scale = min(ph_w / img_w, ph_h / img_h)

    # aplicar tamaño escalado
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)

    pic.width = new_w
    pic.height = new_h

    # centrar dentro del placeholder original
    pic.left = ph_left + (ph_w - new_w) // 2
    pic.top = ph_top + (ph_h - new_h) // 2


def insert_logo_with_scaling(slide, logo_stream):
    """
    Busca un shape que contenga el texto "{{ph_logo}}" e inserta el logo en su lugar,
    manteniendo la relación de aspecto y eliminando el shape original.
    """
    if not logo_stream:
        log("DEBUG: insert_logo_with_scaling - no logo_stream recibido, se omite inserción de logo")
        return False

    found = False
    for shape in slide.shapes:
        if shape.has_text_frame and "{{ph_logo}}" in shape.text:
            log(
                f"DEBUG: insert_logo_with_scaling - encontrado placeholder '{{ph_logo}}' en shape '{shape.name}' con texto '{shape.text.strip()}'"
            )
            if hasattr(logo_stream, "seek"):
                logo_stream.seek(0)
            insert_logo_preserving_aspect(slide, shape, logo_stream)
            found = True
            break  # Solo insertamos el primer logo que encontremos

    if not found:
        log("DEBUG: insert_logo_with_scaling - no se encontró ningún placeholder '{{ph_logo}}' en esta diapositiva")
    return found
