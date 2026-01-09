import os
import json
import base64
import subprocess
import uuid
import warnings
import textwrap
from datetime import datetime, timedelta
from io import BytesIO
from typing import Dict, Any, List

from fastapi import FastAPI, HTTPException, Request
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick
import numpy as np
from pypdf import PdfReader, PdfWriter
from PIL import Image

# Configuración
warnings.filterwarnings("ignore")
DATA_DIR = "/data"
app = FastAPI()

# Modelos de datos
class GenerateRequest(BaseModel):
    data: Dict[str, Any]

MESES_ES = {
    1: "Enero", 2: "Febrero", 3: "Marzo", 4: "Abril",
    5: "Mayo", 6: "Junio", 7: "Julio", 8: "Agosto",
    9: "Septiembre", 10: "Octubre", 11: "Noviembre", 12: "Diciembre"
}

# --------------------------------------------------------------
# UTILS
# --------------------------------------------------------------
def log(msg):
    print(msg, flush=True)

def get_logo_from_base64(base64_string: str) -> BytesIO | None:
    if not base64_string:
        return None
    try:
        image_data = base64.b64decode(base64_string)
        return BytesIO(image_data)
    except Exception:
        log("⚠️ Error de decodificación Base64.")
        return None

def replace_placeholders(slide, replacements):
    for key, value in replacements.items():
        for shape in slide.shapes:
            if shape.has_text_frame:
                texto = shape.text.strip()
                if texto == key:
                    val_str = str(value) if not isinstance(value, str) else value
                    val_str = val_str.replace("\\n", "\n").replace("\\\n", "\n")
                    shape.text = val_str

def insert_image_scaled_by_width(slide, placeholder, image_path_or_stream):
    ph_left, ph_top = placeholder.left, placeholder.top
    ph_width, ph_height = placeholder.width, placeholder.height
    placeholder.element.getparent().remove(placeholder.element)
    pic = slide.shapes.add_picture(image_path_or_stream, ph_left, ph_top, width=ph_width)
    new_height = pic.height
    pic.top = ph_top + (ph_height - new_height) // 2

def insert_logo_preserving_aspect(slide, placeholder, logo_stream):
    ph_left, ph_top = placeholder.left, placeholder.top
    ph_w, ph_h = placeholder.width, placeholder.height
    placeholder.element.getparent().remove(placeholder.element)
    pic = slide.shapes.add_picture(logo_stream, ph_left, ph_top)
    img_w, img_h = pic.width, pic.height
    scale = min(ph_w / img_w, ph_h / img_h)
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)
    pic.width = new_w
    pic.height = new_h
    pic.left = ph_left + (ph_w - new_w) // 2
    pic.top = ph_top + (ph_h - new_h) // 2

def _insert_logo_with_scaling(slide, logo_stream):
    LOGO_PLACEHOLDER_TYPE = 18
    if not logo_stream:
        return
    for shape in slide.placeholders:
        if shape.placeholder_format.type == LOGO_PLACEHOLDER_TYPE:
            insert_logo_preserving_aspect(slide, shape, logo_stream)
            break

def create_composite_logo_from_base64_list(logos_base64_list: list[str], target_height: int = 120) -> BytesIO | None:
    if not logos_base64_list:
        return None

    images = []
    for b64_string in logos_base64_list:
        img_stream = get_logo_from_base64(b64_string)
        if img_stream:
            try:
                img = Image.open(img_stream).convert("RGBA")
                ratio = target_height / img.height
                new_width = int(img.width * ratio)
                img = img.resize((new_width, target_height), Image.LANCZOS)
                images.append(img)
            except Exception as e:
                log(f"⚠️ Error opening image from stream: {e}")

    if not images:
        return None

    max_width = max(img.width for img in images)
    total_height = sum(img.height for img in images)
    composite_image = Image.new("RGBA", (max_width, total_height), (0, 0, 0, 0))

    y_offset = 0
    for img in images:
        x_offset = (max_width - img.width) // 2
        composite_image.paste(img, (x_offset, y_offset), img)
        y_offset += img.height

    output_stream = BytesIO()
    composite_image.save(output_stream, format="PNG")
    output_stream.seek(0)
    return output_stream

def formatea_mes_anio_es(dt: datetime) -> str:
    return f"{MESES_ES.get(dt.month, 'Mes')} {dt.year}"

# --------------------------------------------------------------
# GRÁFICOS
# --------------------------------------------------------------
def create_matplotlib_chart(chart_info, friendly_names, output_file):
    plt.figure(figsize=(10, 5))
    ctype = chart_info.get("type")
    title = chart_info.get("title") or chart_info.get("titulo") or ""
    if title:
        plt.title(title, fontsize=20, fontweight="bold", pad=20)

    plt.xticks(fontsize=18)
    plt.yticks(fontsize=18)
    labels = chart_info.get("labels", [])
    x = range(len(labels))

    flat_friendly_names = {k: v for chart in friendly_names.values() for k, v in chart.items()}
    palette = ["#4f81bd", "#9abb59", "#4bacc6", "#8064a2"]

    series_keys = []
    for key, val in chart_info.items():
        if key in ("labels", "type", "title", "titulo"): continue
        if isinstance(val, (list, tuple)) and all(isinstance(v, (int, float)) for v in val):
            series_keys.append(key)

    if ctype == "bar":
        n = len(series_keys)
        if n > 0:
            ind = np.arange(len(labels))
            total_width = 0.7
            bar_width = total_width / n
            for idx, key in enumerate(series_keys):
                vals = list(chart_info.get(key) or [])
                if len(vals) < len(labels): vals += [0] * (len(labels) - len(vals))
                elif len(vals) > len(labels): vals = vals[:len(labels)]
                
                label_full = flat_friendly_names.get(key, key.replace("_", " ").capitalize())
                label = textwrap.fill(label_full, width=22)
                color = palette[idx % len(palette)]
                offset = (idx - (n - 1) / 2) * bar_width
                plt.bar(ind + offset, vals, bar_width * 0.95, label=label, color=color)
            
            plt.xticks(ind, labels, rotation=0, fontsize=16)
            plt.grid(axis="y", linestyle="-", color="#dcdcdc", linewidth=0.8)
            plt.legend(loc="center left", bbox_to_anchor=(1, 0.5), frameon=False, labelspacing=1.2, fontsize=18)

    elif ctype == "line":
        for idx, key in enumerate(series_keys):
            vals = list(chart_info.get(key) or [])
            if len(vals) < len(labels): vals += [None] * (len(labels) - len(vals))
            elif len(vals) > len(labels): vals = vals[:len(labels)]
            
            label_full = flat_friendly_names.get(key, key.replace("_", " ").capitalize())
            label = textwrap.fill(label_full, width=22)
            color = palette[idx % len(palette)]
            plt.plot(x, vals, label=label, marker="o", color=color)
        
        plt.xticks(x, labels, rotation=45, fontsize=16)
        plt.grid(axis="y", linestyle="-", color="#dcdcdc", linewidth=0.8)
        plt.legend(loc="best", fontsize=18)

    try:
        plt.gca().yaxis.set_major_formatter(mtick.FuncFormatter(lambda x, pos: f"{int(x):,}"))
    except Exception: pass

    plt.tight_layout(rect=[0, 0.03, 0.95, 0.97])
    plt.savefig(output_file, dpi=150, transparent=True)
    plt.close()

def add_charts(slide, charts, friendly_names, replacements_chart):
    to_sort = [s for s in slide.shapes if s.name in replacements_chart]
    chart_placeholders = sorted(to_sort, key=lambda s: s.name)

    for i, (name, chart_info) in enumerate(charts.items()):
        if i >= len(chart_placeholders): break
        placeholder = chart_placeholders[i]
        if not chart_info.get("title") and not chart_info.get("titulo") and name:
            chart_info["title"] = name.replace("_", " ").capitalize()
        
        fn = f"/tmp/{name}_{uuid.uuid4().hex}.png"
        create_matplotlib_chart(chart_info, friendly_names, fn)
        insert_image_scaled_by_width(slide, placeholder, fn)

def chart_builder(values, name, build, kpis):
    total = 0
    for v in values:
        count = sum(values[v])
        total += count
        if count > 0:
            kpis[v] = count

    if total > 0:
        build["charts"][name] = {
            "type": "bar",
            "labels": ["Semana 1", "Semana 2", "Semana 3", "Semana 4"],
            **values,
        }

def build_slide_structure(product, product_name, chart_definitions, pointer_resume):
    slide_info = ["resumen", "sugerencia", "sugerencia_version"]
    build = {
        "titulo": product_name.upper().split(".")[0],
        "kpis": "",
        "charts": {},
    }
    chart_data = {
        chart_name: {serie: [] for serie in series}
        for chart_name, series in chart_definitions.items()
    }
    all_series = {serie: json_key for series in chart_definitions.values() for serie, json_key in series.items()}
    kpis = {}

    for semana in product:
        semana_key = semana.get("Semana", "").strip()
        if semana_key in slide_info:
            valor = semana.get(pointer_resume, "")
            build[semana_key] = valor if valor != "null" else ""
            if semana_key == "sugerencia_version": break
        else:
            for chart_name, series_def in chart_definitions.items():
                for serie_name, json_key in series_def.items():
                    val = semana.get(json_key, 0)
                    try: val = int(float(val))
                    except (ValueError, TypeError): val = 0
                    chart_data[chart_name][serie_name].append(val)

    for chart_name, data in chart_data.items():
        chart_builder(data, chart_name, build, kpis)

    for serie, total in kpis.items():
        nombre_amigable = all_series.get(serie, serie)
        build["kpis"] += f"{nombre_amigable}: {total}\n"

    return build

# --------------------------------------------------------------
# GENERADORES
# --------------------------------------------------------------
# PORTADA
def generar_portada(data, logo_stream):
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_portada.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_titulo}}": data.get("titulo_portada", ""),
        "{{ph_subtitle}}": data.get("subtitulo_portada", ""),
        "{{ph_fecha}}": data.get("fecha_portada", ""),
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    replace_placeholders(slide, replacements)

    # Busca un placeholder de tipo imagen (18) para el logo.
    _insert_logo_with_scaling(slide, logo_stream)

    output = f"{DATA_DIR}/pptx-parts/portada.pptx"
    prs.save(output)
    return output

# CONTENIDO
def generar_contenido(data, logo_stream):
    slides_data = data.get("slides", [])
    generated_files = []

    feet_l, feet_r = data.get("pie_l", ""), data.get("pie_r", "")

    for i, slide_item in enumerate(slides_data):
        template_file = slide_item.get("file_slide", "plantilla_contenido.pptx")
        prs = Presentation(f"{DATA_DIR}/plantillas/{template_file}")
        slide = prs.slides[0]  # Asumimos que la plantilla tiene una sola diapositiva

        slide_content = slide_item.get("slide", {})

        # Cargamos los nombres amigables para las leyendas de los gráficos
        product_type = slide_item.get("type")
        friendly_names = {}
        try:
            with open(
                f"{DATA_DIR}/charts/chart_{product_type}.json", "r", encoding="utf-8"
            ) as f:
                friendly_names = json.load(f)
        except FileNotFoundError:
            log(
                f"⚠️  No se encontró el archivo de configuración de gráficos: chart_{product_type}.json"
            )

        # Diccionario de reemplazos
        replacements = {
            "{{ph_titulo}}": slide_content.get("titulo", ""),
            "{{ph_resumen}}": slide_content.get("resumen", ""),
            "{{ph_sugerencia}}": slide_content.get("sugerencia", ""),
            "{{ph_sugerencia_ver}}": slide_content.get("sugerencia_version", ""),
            "{{ph_pie_l}}": feet_l,
            "{{ph_pie_r}}": feet_r,
        }
        if product_type != "wazuh":
            replacements["{{ph_kpis}}"] = slide_content.get("kpis", "")

        # Reemplazar texto marcador
        replace_placeholders(slide, replacements)
        replacements_chart = [
            "Marcador de posición de imagen 6",
            "Marcador de posición de imagen 9",
            "Marcador de posición de imagen 11",
            "Marcador de posición de imagen 10",
            "Marcador de posición de imagen 12",
        ]
        # Insertar gráficos
        charts = slide_content.get("charts", {})
        if charts:
            add_charts(slide, charts, friendly_names, replacements_chart)

        _insert_logo_with_scaling(slide, logo_stream)

        output_path = f"{DATA_DIR}/pptx-parts/contenido_{product_type}.pptx"
        prs.save(output_path)
        generated_files.append(output_path)
    return generated_files



# CIERRE
def generar_cierre(data, logo_stream):
    cierre = data["despedida"]
    prs = Presentation(f"{DATA_DIR}/plantillas/plantilla_cierre.pptx")
    slide = prs.slides[0]

    replacements = {
        "{{ph_titulo}}": cierre.get("titulo", ""),
        "{{ph_pie_l}}": data.get("pie_l", ""),
        "{{ph_pie_r}}": data.get("pie_r", ""),
    }

    replace_placeholders(slide, replacements)

    _insert_logo_with_scaling(slide, logo_stream)

    output = f"{DATA_DIR}/pptx-parts/cierre.pptx"
    prs.save(output)
    return output


def convert_to_pdf(pptx_file):
    output_dir = f"{DATA_DIR}/pdf-parts"
    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.basename(pptx_file).replace(".pptx", ".pdf")
    pdf_file = os.path.join(output_dir, base_name)

    # Ejecución local de LibreOffice (ya estamos en el contenedor correcto)
    user_inst = f"-env:UserInstallation=file:///tmp/lo_{uuid.uuid4()}"
    cmd = [
        "libreoffice", 
        user_inst, 
        "--headless", 
        "--convert-to", 
        "pdf", 
        pptx_file, 
        "--outdir", 
        output_dir
    ]
    
    try:
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
    except subprocess.CalledProcessError as e:
        log(f"⚠️ Error en LibreOffice: {e.stderr.decode('utf-8', errors='replace')}")
        raise HTTPException(status_code=500, detail=f"LibreOffice error: {e.stderr.decode()}")
    
    return pdf_file

def unir_pdfs(pdf_paths, empresa, type="", split=0):
    writer = PdfWriter()
    for pdf_path in pdf_paths:
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            writer.add_page(page)
    
    output_dir = f"{DATA_DIR}/generados"
    os.makedirs(output_dir, exist_ok=True)
    out_name = f"informe_{empresa}{'.' + type if split == 1 else ''}.pdf"
    out_path = f"{output_dir}/{out_name}"
    
    with open(out_path, "wb") as f:
        writer.write(f)
    return out_path

# --------------------------------------------------------------
# ENDPOINT PRINCIPAL
# --------------------------------------------------------------
@app.post("/generate")
async def generate_report(request: Request):
    try:
        body = await request.json()
        data = body.get("data", {})
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid JSON")

    split = data.get("split", 0)
    logo_stream = get_logo_from_base64(data.get("logo_base64"))

    empresa = data.get("logo")[:-4].lower() if data.get("logo") else ""

    portada = None
    cierre = None
    if data["save"]:
        portada = generar_portada(data, logo_stream)
        cierre = generar_cierre(data, logo_stream)
    contenido_files = generar_contenido(data, logo_stream)
    types = [slide.get("type", "") for slide in data.get("slides", [])]

    # Pre-convert common parts to avoid redundant processing
    portada_pdf = convert_to_pdf(portada) if (data["save"] and portada) else None
    cierre_pdf = convert_to_pdf(cierre) if (data["save"] and cierre) else None

    informe_names = []
    if split == 0:
        pdf_files_to_merge = []
        if portada_pdf:
            pdf_files_to_merge.append(portada_pdf)
        
        pdf_files_to_merge.extend([convert_to_pdf(f) for f in contenido_files])
        
        if cierre_pdf:
            pdf_files_to_merge.append(cierre_pdf)

        informe_names.append(unir_pdfs(pdf_files_to_merge, empresa))
    else:
        for idx, content_pptx in enumerate(contenido_files):
            
            pdf_files_to_merge = []
            if portada_pdf:
                pdf_files_to_merge.append(portada_pdf)
            
            pdf_files_to_merge.append(convert_to_pdf(content_pptx))
            
            if cierre_pdf:
                pdf_files_to_merge.append(cierre_pdf)

            informe_names.append(
                unir_pdfs(pdf_files_to_merge, empresa, types[idx], split)
            )

    return {"file_names": informe_names}

@app.post("/build-structure")
async def build_structure(request: Request):
    try:
        body = await request.json()
        data = body.get("data", {})
        main = data.get("main", {})
        products = data.get("products", [])
        
        # Procesar fecha portada
        fecha_iso = main.get("fecha_portada")
        if fecha_iso:
            try:
                dt_object = None
                if isinstance(fecha_iso, str):
                    try: dt_object = datetime.strptime(fecha_iso, "%Y-%m")
                    except ValueError:
                        try: dt_object = datetime.fromisoformat(fecha_iso.replace("Z", "+00:00"))
                        except Exception: dt_object = None
                    if dt_object: dt_object = dt_object.replace(day=1)
                elif isinstance(fecha_iso, (int, float)):
                    dt_object = datetime(1899, 12, 30) + timedelta(days=fecha_iso)
                    dt_object = dt_object.replace(day=1)
                
                if dt_object: main["fecha_portada"] = formatea_mes_anio_es(dt_object)
                else: main["fecha_portada"] = "Fecha no válida"
            except Exception: main["fecha_portada"] = "Fecha no válida"
        else:
            main["fecha_portada"] = "Fecha no válida"

        # Agrupar productos
        parse_products = {}
        actual_product = ""
        for product in products:
            actual_product = product.get("product", actual_product)
            if actual_product not in parse_products:
                parse_products[actual_product] = []
            parse_products[actual_product].append(product)

        # Construir slides
        if "slides" not in main: main["slides"] = []
        
        file_slide_map = {
            "uas": "plantilla_contenido.pptx",
            "wazuh": "plantilla_contenido_no_kpis.pptx",
            "ardid": "plantilla_contenido.pptx",
            "invgate.asj": "plantilla_contenido_no_kpis.pptx",
            "invgate": "plantilla_contenido.pptx",
            "beyondtrust": "plantilla_contenido.pptx",
            "whalemate": "plantilla_contenido.pptx",
        }

        for product_key in parse_products:
            if not parse_products[product_key]: continue
            pointer_resumen = list(parse_products[product_key][0].keys())[1]
            
            chart_def = {}
            try:
                with open(f"{DATA_DIR}/charts/chart_{product_key}.json", "r", encoding="utf-8") as f:
                    chart_def = json.load(f)
            except FileNotFoundError: pass

            slide_data = build_slide_structure(parse_products[product_key], product_key, chart_def, pointer_resumen)
            if slide_data:
                main["slides"].append({
                    "type": product_key,
                    "slide": slide_data,
                    "file_slide": file_slide_map.get(product_key, "plantilla_contenido.pptx")
                })

        return {"status": "ok", "output_file": main}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error building structure: {str(e)}")

@app.post("/generate-n-emp")
async def generate_pdf_n_emp(request: Request):
    try:
        body = await request.json()
        data = body.get("data", {})
        main_data = data.get("main", {})
        emp_codes = data.get("emp_codes", [])
        logos_base64_list = data.get("logos_base64", [])

        logo_stream = create_composite_logo_from_base64_list(logos_base64_list)
        empresa = "-".join(emp_codes).lower()

        # Generar Portada
        portada_path = generar_portada(main_data, logo_stream)
        portada_pdf = convert_to_pdf(portada_path)

        # Generar Cierre
        cierre_path = generar_cierre(main_data, logo_stream)
        cierre_pdf = convert_to_pdf(cierre_path)

        # Unir todo
        pdf_files_to_merge = [portada_pdf]
        pdf_files_to_merge.extend([os.path.join(DATA_DIR, "generados", f"informe_{f.lower()}.pdf") for f in emp_codes])
        pdf_files_to_merge.append(cierre_pdf)
        
        final_pdf = unir_pdfs(pdf_files_to_merge, empresa)
        return {"file_name": os.path.basename(f"informe_{empresa}")}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating PDF N Emp: {str(e)}")

@app.get("/health")
def health():
    return {"status": "ok"}